# -*- coding: utf-8 -*-
"""
build_pptx.py — 사업계획서 PPTX 렌더러 (맥킨지 파스텔)

사용:
    python build_pptx.py <slides.json> <output.pptx>

디자인 규칙 (pptx-pastel-style 계승):
- 파스텔 5색 순환(하늘·연주황·보라·노랑·민트) + 네이비 앵커
- 1 슬라이드 1 메시지 — 모든 content 슬라이드에 액션 타이틀
- Bullet ≤ 4 (초과 시 경고 + 종료 코드 3)
- 전 슬라이드 푸터: "내부 검토용 — 투자권유 아님" + 페이지 번호
"""
import sys
import json

sys.stdout.reconfigure(encoding="utf-8", errors="replace")
sys.stderr.reconfigure(encoding="utf-8", errors="replace")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("[ERROR] python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

# ---------------- 팔레트 (html-pastel-style 토큰 계승)
NAVY = RGBColor(0x16, 0x38, 0x6E)
NAVY_BG = RGBColor(0xEA, 0xF1, 0xFB)
INK = RGBColor(0x19, 0x22, 0x3C)
TEXT = RGBColor(0x37, 0x41, 0x5C)
MUTED = RGBColor(0x6C, 0x74, 0x8C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FAMS = [  # (bg, accent, deep)
    (RGBColor(0xEA, 0xF2, 0xFE), RGBColor(0x3D, 0x7B, 0xEA), RGBColor(0x1B, 0x4D, 0xA0)),  # sky
    (RGBColor(0xFC, 0xEB, 0xDA), RGBColor(0xEA, 0x6A, 0x17), RGBColor(0xA8, 0x50, 0x0F)),  # orange
    (RGBColor(0xF1, 0xEA, 0xFB), RGBColor(0x87, 0x5A, 0xD2), RGBColor(0x65, 0x3A, 0x9F)),  # violet
    (RGBColor(0xFB, 0xF3, 0xD9), RGBColor(0xD2, 0xA1, 0x1C), RGBColor(0x8A, 0x69, 0x12)),  # amber
    (RGBColor(0xE6, 0xF5, 0xEE), RGBColor(0x2C, 0x9F, 0x70), RGBColor(0x1C, 0x74, 0x4E)),  # mint
]
FONT = "Pretendard"
SW, SH = Inches(13.333), Inches(7.5)  # 16:9
WARNINGS = []


def warn(msg):
    WARNINGS.append(msg)
    print(f"[WARN] {msg}", file=sys.stderr)


def box(slide, x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.adjustments[0] = 0.06
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill if fill else WHITE
    if line:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def put_text(shape_or_slide, x, y, w, h, runs, size=14, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    """runs: str 또는 [(text, {bold,color,size}), ...] 목록의 목록(문단별)"""
    tb = shape_or_slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        if isinstance(para, str):
            para = [(para, {})]
        for txt, st in para:
            r = p.add_run()
            r.text = txt
            r.font.name = FONT
            r.font.size = Pt(st.get("size", size))
            r.font.bold = st.get("bold", bold)
            r.font.color.rgb = st.get("color", color)
    return tb


def footer(slide, page_no, total):
    put_text(slide, Inches(0.5), SH - Inches(0.42), Inches(6), Inches(0.3),
             "내부 검토용 — 투자권유 아님 · 외부 배포 금지", size=9, color=MUTED)
    put_text(slide, SW - Inches(1.2), SH - Inches(0.42), Inches(0.7), Inches(0.3),
             f"{page_no} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def fam_of(idx):
    return FAMS[idx % len(FAMS)]


# ---------------- 레이아웃별 렌더
def render_cover(prs, meta):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box(s, 0, 0, SW, SH, fill=NAVY_BG)
    box(s, Inches(0.7), Inches(1.1), SW - Inches(1.4), Inches(4.2), fill=WHITE)
    put_text(s, Inches(1.1), Inches(1.5), SW - Inches(2.2), Inches(0.4),
             meta.get("eyebrow", "신사업 사업계획서 — 내부 투자심의용"),
             size=13, color=MUTED, bold=True)
    put_text(s, Inches(1.1), Inches(2.0), SW - Inches(2.2), Inches(1.4),
             meta.get("title", "사업계획서"), size=40, color=NAVY, bold=True)
    put_text(s, Inches(1.1), Inches(3.4), SW - Inches(2.2), Inches(0.6),
             meta.get("subtitle", ""), size=18, color=TEXT)
    put_text(s, Inches(1.1), Inches(4.3), SW - Inches(2.2), Inches(0.5),
             f"{meta.get('org', '')}   ·   {meta.get('date', '')}", size=13, color=MUTED)
    # 면책 배지 (헌법 4조 — 자동 삽입, json에 없어도 강제)
    b = box(s, Inches(0.7), Inches(5.7), SW - Inches(1.4), Inches(0.9),
            fill=RGBColor(0xFD, 0xEE, 0xE6), line=RGBColor(0xF8, 0xD6, 0xC5))
    put_text(s, Inches(1.0), Inches(5.85), SW - Inches(2.0), Inches(0.6),
             "본 자료는 내부 의사결정 검토용이며 투자권유가 아닙니다. "
             "수치에는 [가정] 추정이 포함되어 있어 외부 배포를 금합니다.",
             size=11, color=RGBColor(0xB0, 0x51, 0x2B))
    return s


def render_section(prs, sl, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg, ac, dp = fam_of(idx)
    box(s, 0, 0, SW, SH, fill=WHITE)
    box(s, 0, Inches(2.9), SW, Inches(1.7), fill=bg)
    put_text(s, Inches(0.9), Inches(3.1), Inches(1.6), Inches(0.6),
             sl.get("badge", ""), size=16, color=ac, bold=True)
    put_text(s, Inches(0.9), Inches(3.55), SW - Inches(1.8), Inches(0.9),
             sl.get("title", ""), size=28, color=dp, bold=True)
    return s


def render_content(prs, sl, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg, ac, dp = fam_of(idx)
    box(s, 0, 0, SW, Inches(0.16), fill=ac)  # 상단 액센트 바
    put_text(s, Inches(0.6), Inches(0.35), SW - Inches(1.2), Inches(0.45),
             sl.get("title", ""), size=15, color=MUTED, bold=True)
    put_text(s, Inches(0.6), Inches(0.82), SW - Inches(1.2), Inches(0.75),
             sl.get("action_title", ""), size=21, color=INK, bold=True)
    bullets = sl.get("bullets", [])
    if len(bullets) > 4:
        warn(f"슬라이드 '{sl.get('title')}' bullet {len(bullets)}개 — 4개 이하 규칙 위반")
    y = Inches(1.8)
    card = box(s, Inches(0.6), y, SW - Inches(1.2), SH - y - Inches(0.75), fill=bg)
    by = y + Inches(0.35)
    for b in bullets[:6]:
        put_text(s, Inches(1.0), by, SW - Inches(2.0), Inches(0.75),
                 [[("●  ", {"color": ac, "size": 12}), (str(b), {})]],
                 size=15, color=TEXT)
        by += Inches(0.85)
    if sl.get("note"):
        put_text(s, Inches(1.0), SH - Inches(1.35), SW - Inches(2.0), Inches(0.5),
                 sl["note"], size=11.5, color=MUTED)
    return s


def render_table(prs, sl, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg, ac, dp = fam_of(idx)
    box(s, 0, 0, SW, Inches(0.16), fill=ac)
    put_text(s, Inches(0.6), Inches(0.35), SW - Inches(1.2), Inches(0.45),
             sl.get("title", ""), size=15, color=MUTED, bold=True)
    put_text(s, Inches(0.6), Inches(0.82), SW - Inches(1.2), Inches(0.75),
             sl.get("action_title", ""), size=21, color=INK, bold=True)
    t = sl.get("table", {})
    headers = t.get("headers", [])
    rows = t.get("rows", [])
    if not headers:
        warn(f"table 슬라이드 '{sl.get('title')}' headers 없음 — 건너뜀")
        return s
    n_r, n_c = len(rows) + 1, len(headers)
    top = Inches(1.85)
    height = min(SH - top - Inches(0.85), Inches(0.42) * n_r)
    shp = s.shapes.add_table(n_r, n_c, Inches(0.6), top, SW - Inches(1.2), height)
    tbl = shp.table
    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = str(htxt)
        c.fill.solid()
        c.fill.fore_color.rgb = ac
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True
                r.font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        for j in range(n_c):
            c = tbl.cell(i, j)
            c.text = str(row[j]) if j < len(row) else ""
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else bg
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = FONT; r.font.size = Pt(11.5)
                    r.font.color.rgb = TEXT
    if sl.get("note"):
        put_text(s, Inches(0.6), SH - Inches(0.85), SW - Inches(1.2), Inches(0.4),
                 sl["note"], size=11, color=MUTED)
    return s


def render_kpi(prs, sl, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg, ac, dp = fam_of(idx)
    box(s, 0, 0, SW, Inches(0.16), fill=ac)
    put_text(s, Inches(0.6), Inches(0.35), SW - Inches(1.2), Inches(0.45),
             sl.get("title", ""), size=15, color=MUTED, bold=True)
    put_text(s, Inches(0.6), Inches(0.82), SW - Inches(1.2), Inches(0.75),
             sl.get("action_title", ""), size=21, color=INK, bold=True)
    kpis = sl.get("kpis", [])[:4]
    if not kpis:
        warn(f"kpi 슬라이드 '{sl.get('title')}' kpis 없음")
        return s
    n = len(kpis)
    gap = Inches(0.25)
    total_w = SW - Inches(1.2)
    w = Emu(int((total_w - gap * (n - 1)) / n))
    x = Inches(0.6)
    for i, k in enumerate(kpis):
        fbg, fac, fdp = fam_of(idx + i)
        box(s, x, Inches(2.2), w, Inches(2.6), fill=fbg)
        put_text(s, x + Inches(0.25), Inches(2.5), w - Inches(0.5), Inches(0.5),
                 str(k.get("label", "")), size=12.5, color=fdp, bold=True)
        put_text(s, x + Inches(0.25), Inches(3.05), w - Inches(0.5), Inches(0.9),
                 str(k.get("value", "")), size=27, color=fdp, bold=True)
        put_text(s, x + Inches(0.25), Inches(4.0), w - Inches(0.5), Inches(0.7),
                 str(k.get("desc", "")), size=11, color=TEXT)
        x = Emu(int(x + w + gap))
    if sl.get("note"):
        put_text(s, Inches(0.6), SH - Inches(1.5), SW - Inches(1.2), Inches(0.5),
                 sl["note"], size=11.5, color=MUTED)
    return s


RENDER = {"cover": None, "section": render_section, "content": render_content,
          "table": render_table, "kpi": render_kpi}


def main():
    argv = [a for a in sys.argv[1:] if not a.startswith("--")]
    allow_warn = "--allow-warnings" in sys.argv
    if len(argv) < 2:
        print("Usage: python build_pptx.py <slides.json> <output.pptx> [--allow-warnings]",
              file=sys.stderr)
        sys.exit(1)
    cfg_path, out_path = argv[0], argv[1]
    if not out_path.lower().endswith(".pptx"):
        print(f"[ERROR] 출력 파일은 .pptx 여야 합니다: {out_path}", file=sys.stderr)
        sys.exit(1)
    try:
        with open(cfg_path, "r", encoding="utf-8") as f:
            cfg = json.load(f)
    except FileNotFoundError:
        print(f"[ERROR] 입력 json 없음: {cfg_path}", file=sys.stderr)
        sys.exit(1)
    except json.JSONDecodeError as e:
        print(f"[ERROR] json 파싱 실패: {e}", file=sys.stderr)
        sys.exit(1)
    meta = cfg.get("meta")
    slides = cfg.get("slides")
    if not meta or not isinstance(slides, list) or not slides:
        print("[ERROR] config에 meta 와 slides[] 가 필요합니다. slides.sample.json 참조.",
              file=sys.stderr)
        sys.exit(1)

    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH

    render_cover(prs, meta)
    sec_idx = 0
    for sl in slides:
        layout = sl.get("layout", "content")
        if layout == "section":
            fn = render_section
            fn(prs, sl, sec_idx)
            sec_idx += 1
        else:
            fn = RENDER.get(layout)
            if fn is None:
                warn(f"알 수 없는 layout '{layout}' — content로 대체")
                fn = render_content
            fn(prs, sl, max(0, sec_idx - 1))

    total = len(prs.slides._sldIdLst)
    for i, s in enumerate(prs.slides, start=1):
        if i == 1:
            continue  # 표지는 푸터 생략
        footer(s, i, total)

    prs.save(out_path)

    # 자체 검증 — 다시 열어 슬라이드 수 확인
    chk = Presentation(out_path)
    n = len(chk.slides._sldIdLst)
    print(f"OK: {out_path}")
    print(f"slides={n}")
    print(f"warnings={len(WARNINGS)}")
    if n != len(slides) + 1:
        print(f"[ERROR] 슬라이드 수 불일치: 기대 {len(slides)+1}, 실제 {n}", file=sys.stderr)
        sys.exit(2)
    if WARNINGS and not allow_warn:
        print("[FAIL] 경고가 있어 실패로 처리합니다 (bullet>4 등). 수정하거나 "
              "--allow-warnings 를 붙이고 보고서에 사유를 명시하십시오.", file=sys.stderr)
        sys.exit(3)


if __name__ == "__main__":
    main()
